#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI热点监控工具 - PPT生成脚本
需要安装: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """创建AI热点监控工具项目演示PPT"""

    # 创建演示文稿对象
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色定义
    PRIMARY_COLOR = RGBColor(102, 126, 234)  # #667eea
    SECONDARY_COLOR = RGBColor(118, 75, 162)  # #764ba2
    ACCENT_COLOR = RGBColor(240, 147, 251)    # #f093fb
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    GRAY = RGBColor(51, 51, 51)

    def add_title_slide(title, subtitle):
        """添加标题页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 设置渐变背景（模拟）
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135
        fill.gradient_stops[0].color.rgb = PRIMARY_COLOR
        fill.gradient_stops[1].color.rgb = SECONDARY_COLOR

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE

        # 添加副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2), Inches(11.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = WHITE

    def add_content_slide(title, content_items, is_two_column=False):
        """添加内容页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 白色背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # 添加装饰线
        line = slide.shapes.add_shape(
            1,  # 直线
            Inches(0.5), Inches(1.2), Inches(12.333), Inches(0)
        )
        line.line.color.rgb = SECONDARY_COLOR
        line.line.width = Pt(4)

        # 添加内容
        if is_two_column:
            # 两列布局
            left = Inches(0.7)
            top = Inches(1.6)
            width = Inches(5.8)

            for i, item in enumerate(content_items):
                col = i % 2
                row = i // 2
                content_box = slide.shapes.add_textbox(
                    left + col * 6, top + row * 0.7, width, Inches(0.6)
                )
                content_frame = content_box.text_frame
                content_frame.text = item
                content_para = content_frame.paragraphs[0]
                content_para.font.size = Pt(20)
                content_para.font.color.rgb = GRAY
                content_para.space_before = Pt(10)
        else:
            # 单列布局
            content_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.6), Inches(12), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, item in enumerate(content_items):
                if i > 0:
                    content_frame.add_paragraph()
                p = content_frame.paragraphs[i]
                p.text = item
                p.font.size = Pt(22)
                p.font.color.rgb = GRAY
                p.space_before = Pt(15)
                p.level = 0

    def add_cards_slide(title, cards):
        """添加卡片布局页面"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 白色背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # 添加卡片
        for i, card in enumerate(cards):
            row = i // 2
            col = i % 2
            left = Inches(0.7 + col * 6.2)
            top = Inches(1.6 + row * 2.6)

            # 卡片背景
            card_shape = slide.shapes.add_shape(
                1,  # 矩形
                left, top, Inches(5.8), Inches(2.3)
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
            card_shape.line.color.rgb = PRIMARY_COLOR
            card_shape.line.width = Pt(2)

            # 卡片标题
            card_title = slide.shapes.add_textbox(
                left + Inches(0.2), top + Inches(0.2), Inches(5.4), Inches(0.5)
            )
            card_title_frame = card_title.text_frame
            card_title_frame.text = card['title']
            card_title_para = card_title_frame.paragraphs[0]
            card_title_para.font.size = Pt(24)
            card_title_para.font.bold = True
            card_title_para.font.color.rgb = PRIMARY_COLOR

            # 卡片内容
            card_content = slide.shapes.add_textbox(
                left + Inches(0.2), top + Inches(0.7), Inches(5.4), Inches(1.4)
            )
            card_content_frame = card_content.text_frame
            card_content_frame.word_wrap = True
            card_content_frame.text = card['content']
            for para in card_content_frame.paragraphs:
                para.font.size = Pt(18)
                para.font.color.rgb = GRAY

    # 第1页：封面
    add_title_slide(
        "AI热点监控工具",
        "智能化的热点发现和通知系统\n\n7×24小时热点监控 · AI内容验证 · 多渠道通知"
    )

    # 第2页：目录
    add_content_slide(
        "目录",
        [
            "项目背景与痛点",
            "核心功能介绍",
            "技术架构设计",
            "项目亮点特色",
            "技术栈展示",
            "功能演示",
            "未来规划"
        ]
    )

    # 第3页：项目背景
    add_cards_slide(
        "项目背景 - 痛点问题",
        [
            {"title": "效率低下", "content": "人工搜索热点需要浏览多个网站，耗时耗力"},
            {"title": "容易遗漏", "content": "无法实现7×24小时持续监控，错过重要信息"},
            {"title": "虚假内容", "content": "AI编程领域虚假信息多，难以辨别真伪"},
            {"title": "通知延迟", "content": "无法第一时间获取热点，错失最佳时机"}
        ]
    )

    # 第4页：核心功能（一）
    add_cards_slide(
        "核心功能（一）",
        [
            {
                "title": "智能热点发现",
                "content": "• 自动监控多个科技新闻网站\n• 基于关键词智能匹配\n• 实时内容分析和去重\n• 支持自定义搜索范围"
            },
            {
                "title": "AI内容验证",
                "content": "• 集成OpenRouter API\n• 智能识别虚假内容\n• 提供可信度评分\n• 生成详细分析报告"
            }
        ]
    )

    # 第5页：核心功能（二）
    add_cards_slide(
        "核心功能（二）",
        [
            {
                "title": "多渠道通知",
                "content": "• 邮件通知\n• Web推送通知\n• Slack集成\n• 通知历史记录"
            },
            {
                "title": "关键词管理",
                "content": "• 添加/编辑/删除关键词\n• 关键词分组管理\n• 优先级设置\n• 实时监控状态"
            }
        ]
    )

    # 第6页：技术架构
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 白色背景
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "技术架构"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # 架构图文本框
    arch_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(11.333), Inches(5)
    )
    arch_frame = arch_box.text_frame
    arch_frame.text = (
        "┌─────────────────────────────────────────┐\n"
        "│         前端界面 (React)                 │\n"
        "├─────────────────────────────────────────┤\n"
        "│         API层 (Express)                  │\n"
        "├─────────────────────────────────────────┤\n"
        "│  业务逻辑层 (Services + Controllers)     │\n"
        "├─────────────────────────────────────────┤\n"
        "│ AI服务  爬虫服务  通知服务               │\n"
        "├─────────────────────────────────────────┤\n"
        "│         数据层 (MongoDB)                 │\n"
        "└─────────────────────────────────────────┘"
    )
    for para in arch_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.name = "Courier New"
        para.font.color.rgb = BLACK

    # 第7页：技术栈
    add_content_slide(
        "技术栈",
        [
            "后端技术: Node.js 18+ | Express.js | MongoDB | JWT | node-cron | axios + cheerio",
            "前端技术: React 18+ | TypeScript | Ant Design | React Router | Recharts",
            "AI服务: OpenRouter API | DeepSeek API",
            "开发工具: Git | Docker | npm | TypeScript ESLint",
            "部署方案: Docker容器化 | 传统Node.js部署 | 环境变量配置"
        ]
    )

    # 第8页：项目亮点
    add_cards_slide(
        "项目亮点",
        [
            {
                "title": "双AI服务商支持",
                "content": "OpenRouter: GPT-4、Claude等\nDeepSeek: 国内稳定服务\n环境变量一键切换"
            },
            {
                "title": "独特视觉设计",
                "content": "紫粉渐变现代设计\n卡片式布局\n完美响应式适配"
            },
            {
                "title": "Agent Skills封装",
                "content": "标准化接口设计\n支持第三方AI集成\n三大核心Skill"
            },
            {
                "title": "完善工程实践",
                "content": "TypeScript类型安全\nDocker容器化\n完整项目文档"
            }
        ]
    )

    # 第9页：功能演示
    add_content_slide(
        "功能演示",
        [
            "Dashboard仪表板 - 统计卡片、热点列表、渐变设计风格",
            "关键词管理 - 卡片式布局、搜索过滤、状态管理",
            "通知中心 - 列表展示、批量操作、详情查看",
            "系统设置 - 通知配置、扫描设置、账户管理"
        ],
        is_two_column=True
    )

    # 第10页：未来规划
    add_content_slide(
        "未来规划",
        [
            "短期目标: 前后端集成测试 | 性能优化 | 安全加固 | 监控告警",
            "长期规划: 更多AI模型支持 | 扩展数据源 | 移动端App | 国际化支持",
            "项目统计: 30+后端文件 | 10+前端文件 | 3个核心Skills"
        ]
    )

    # 第11页：总结
    add_content_slide(
        "项目总结",
        [
            "项目价值: 自动化热点发现 | AI内容验证 | 7×24小时监控 | 灵活架构设计",
            "核心优势: 双AI服务商支持 | 独特视觉设计 | 完整Skills封装 | 完善工程实践",
            "",
            "让AI热点监控帮助您走在科技前沿！"
        ]
    )

    # 第12页：致谢页
    add_title_slide(
        "谢谢观看！",
        "欢迎提问和讨论\n\n📧 项目地址: github.com/your-repo\n🐛 问题反馈: github.com/your-repo/issues"
    )

    # 保存演示文稿
    output_file = "AI热点监控工具_项目分享.pptx"
    prs.save(output_file)
    print(f"✓ PPT已生成: {output_file}")
    print(f"✓ 共 {len(prs.slides)} 页")

if __name__ == "__main__":
    print("正在生成AI热点监控工具项目演示PPT...")
    try:
        create_presentation()
        print("\n生成完成！")
    except ImportError:
        print("错误: 需要先安装 python-pptx 库")
        print("请运行: pip install python-pptx")
    except Exception as e:
        print(f"错误: {e}")
